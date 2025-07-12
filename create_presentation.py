from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_content_moderation_presentation():
    """Create a PowerPoint presentation for the Content Moderation ADK project."""
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Objective and Problem
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
    title1 = slide1.shapes.title
    title1.text = "🎯 Objective & Problem Statement"
    
    content1 = slide1.placeholders[1]
    content1.text = """🔍 PROBLEM STATEMENT:
• Manual content moderation is time-consuming and error-prone
• Need for automated detection of inappropriate content
• Multiple violation types require different detection approaches
• Lack of comprehensive confidence scoring system

🎯 OBJECTIVES:
• Develop automated content moderation pipeline
• Detect multiple violation types: nudity, violence, drugs, alcohol, hate symbols, PII, QR codes
• Implement confidence-based decision making
• Create scalable and efficient moderation system
• Provide detailed reporting and analysis

📊 CHALLENGES:
• Balancing accuracy vs. speed
• Handling edge cases and exceptions
• Managing false positives/negatives
• Integrating multiple AI models effectively"""

    # Slide 2: Solution
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "🛡️ Solution Overview"
    
    content2 = slide2.placeholders[1]
    content2.text = """🔧 CORE SOLUTION:
• Multi-agent content moderation pipeline
• Hybrid approach: NudeNet + Gemini Vision AI
• Confidence-based decision making system
• Comprehensive violation detection

🧩 AGENT COMPONENTS:
• Ingestion Agent: Image preprocessing and validation
• Nudity Detection: NudeNet for accurate nudity detection
• Nudity Exceptions: Gemini for context-aware exceptions
• Violence Detection: Gemini Vision for blood, weapons, gore
• Drugs Detection: Gemini Vision for paraphernalia
• Alcohol/Smoking: Gemini Vision for substance detection
• Hate Symbols: Gemini Vision for extremist content
• PII/Text Detection: Gemini Vision for personal information
• QR Code Detection: Gemini Vision for embedded codes

⚡ KEY FEATURES:
• Real-time processing capabilities
• Detailed confidence scoring
• Comprehensive reporting system
• Scalable architecture
• Exception handling and edge case management"""

    # Slide 3: Architecture Overview
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "🏗️ Architecture Overview"
    
    content3 = slide3.placeholders[1]
    content3.text = """📋 PIPELINE ARCHITECTURE:

1️⃣ INGESTION LAYER:
   • Image preprocessing and validation
   • Format conversion and optimization
   • Quality assessment

2️⃣ DETECTION LAYER:
   • NudeNet: Primary nudity detection
   • Gemini Vision: Context-aware analysis
   • Multi-modal violation detection

3️⃣ ANALYSIS LAYER:
   • Confidence score calculation
   • Violation type classification
   • Decision logic processing

4️⃣ OUTPUT LAYER:
   • Final decision (Accept/Reject/Flag)
   • Detailed reporting
   • JSON export capabilities

🔄 PROCESS FLOW:
Image Input → Preprocessing → Multi-Agent Detection → 
Confidence Analysis → Decision Making → Report Generation

📊 CONFIDENCE SYSTEM:
• 0.9: Very High (Strong evidence)
• 0.8: High (Reliable detection)
• 0.6: Medium-High (Good evidence)
• 0.5: Medium (Needs review)
• 0.4: Medium-Low (Weak evidence)
• 0.1: Low (Safe content)
• 0.0: Error (Manual review)"""

    # Slide 4: Pipeline Diagram
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout
    title4 = slide4.shapes.title
    title4.text = "🔄 Moderation Pipeline Flow"

    # Draw blocks and arrows
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(1.6)
    height = Inches(0.8)
    gap = Inches(0.3)

    stages = [
        "Image Input",
        "Ingestion",
        "Multi-Agent Detection\n(Nudity, Violence, Drugs, etc.)",
        "Confidence Analysis",
        "Decision",
        "Report"
    ]

    shapes = slide4.shapes
    block_shapes = []
    for i, stage in enumerate(stages):
        block = shapes.add_shape(
            1,  # Rectangle
            left + i * (width + gap),
            top,
            width,
            height
        )
        block.text = stage
        block.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        block.text_frame.paragraphs[0].font.size = Pt(14)
        block.text_frame.paragraphs[0].font.bold = True
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor(0x4F, 0x81, 0xBD)
        block.line.color.rgb = RGBColor(0, 0, 0)
        block_shapes.append(block)
        # Add arrow except after last block
        if i < len(stages) - 1:
            arrow = shapes.add_connector(
                2,  # Arrow
                left + (i + 1) * (width + gap) - gap,
                top + height / 2 - 0.05,
                gap,
                0.1
            )
            arrow.line.color.rgb = RGBColor(0, 0, 0)

    # Slide 5: Efficiency Metrics
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "📈 Efficiency Metrics"
    
    content5 = slide5.placeholders[1]
    content5.text = """🎯 PERFORMANCE METRICS:

⚡ SPEED & EFFICIENCY:
• Real-time processing: < 30 seconds per image
• Parallel agent execution capability
• Optimized image preprocessing
• Efficient API integration

🎯 ACCURACY METRICS:
• NudeNet: Industry-standard nudity detection
• Gemini Vision: Advanced context understanding
• Confidence-based decision making
• Reduced false positive/negative rates

📊 DETECTION CAPABILITIES:
• Nudity: 95%+ accuracy with NudeNet
• Violence: Multi-category detection (blood, weapons, gore)
• Drugs: Paraphernalia and substance detection
• Alcohol/Smoking: Context-aware substance detection
• Hate Symbols: Symbol and extremist content detection
• PII: Personal information and text extraction
• QR Codes: Embedded code detection

🔄 SCALABILITY:
• Modular agent architecture
• Easy addition of new detection types
• Configurable confidence thresholds
• Exportable results for analysis

📋 REPORTING FEATURES:
• Comprehensive JSON reports
• Confidence score breakdown
• Violation type classification
• Decision rationale explanation"""

    # Save the presentation
    output_path = "Content_Moderation_ADK_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    print(f"📁 Location: {os.path.abspath(output_path)}")
    
    return output_path

if __name__ == "__main__":
    create_content_moderation_presentation() 